"""
generate-ppt.py
Serverless function: accepts POST with JSON body, downloads a master PPTX
template, fills placeholders, and returns a binary .pptx download.

Expected JSON body:
{
    "template_url": "https://...",
    "date": "Sunday, 29 June 2026",
    "sections": [
        {
            "name": "Entrance",
            "song": {
                "title": "Song Title",
                "lyrics": [
                    { "label": "Verse 1", "text": "Line 1\\nLine 2" },
                    { "label": "Chorus",  "text": "Chorus line 1\\nChorus line 2" },
                    { "label": "Verse 2", "text": "Line 3\\nLine 4" }
                ]
            }
        },
        ... (9 sections total)
    ]
}

Slide mapping (0-indexed):
  0  = Cover Page     (replace {{DATE}})
  1  = Entrance
  2  = Lord Have Mercy
  3  = Gloria
  4  = Acclamation
  5  = Offertory
  6  = Holy Holy
  7  = Proclamation
  8  = Communion
  9  = Recessional
"""

import copy
import json
import io
import re
import traceback
from http.server import BaseHTTPRequestHandler
from urllib.parse import urlparse, parse_qs, unquote

import requests
from pptx import Presentation
from pptx.util import Pt, Emu
from lxml import etree
from pypdf import PdfWriter, PdfReader


# ---------------------------------------------------------------------------
# OOXML DrawingML namespace
# ---------------------------------------------------------------------------
_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Section order (must match slide indices 1-9 in the template)
# ---------------------------------------------------------------------------
SECTION_ORDER = [
    "Entrance",
    "Lord Have Mercy",
    "Gloria",
    "Acclamation",
    "Offertory",
    "Holy Holy",
    "Proclamation",
    "Communion",
    "Recessional",
]


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def _normalise(text: str) -> str:
    """
    Convert literal backslash-n sequences (as stored in some databases)
    into real newline characters, and strip surrounding whitespace.
    """
    return text.replace("\\n", "\n").strip()


def _para_full_text(para) -> str:
    """Concatenate all run texts in a paragraph (handles split-run placeholders)."""
    return "".join(run.text for run in para.runs)


def _frame_contains(text_frame, placeholder: str) -> bool:
    """
    Return True if *placeholder* exists anywhere in the text frame.
    Checks concatenated paragraph text so split-run placeholders are found.
    """
    for para in text_frame.paragraphs:
        if placeholder in _para_full_text(para):
            return True
    return False


def _replace_in_para(para, old: str, new: str) -> bool:
    """
    Replace *old* with *new* within a paragraph, even if the placeholder
    is split across multiple runs (consolidates into the first run).
    """
    full = _para_full_text(para)
    if old not in full:
        return False
    replaced = full.replace(old, new)
    runs = para.runs
    if runs:
        runs[0].text = replaced
        for run in runs[1:]:
            run.text = ""
    return True


def _replace_text_in_frame(text_frame, old: str, new: str) -> bool:
    """Replace *old* with *new* across every paragraph in a text frame."""
    changed = False
    for para in text_frame.paragraphs:
        if _replace_in_para(para, old, new):
            changed = True
    return changed


def _snapshot_rPr(text_frame):
    """
    Return a deep-copied <a:rPr> element from the first run in the frame,
    or None if the frame has no runs with run properties.
    Used to carry over font name, size, colour when building new paragraphs.
    """
    for para in text_frame.paragraphs:
        for run in para.runs:
            rPr_el = run._r.find(f"{{{_NS}}}rPr")
            if rPr_el is not None:
                return copy.deepcopy(rPr_el)
    return None


def _parse_lines_with_bold(text: str, base_bold: bool = False) -> list[tuple[str, bool]]:
    """
    Split *text* by ``**`` markers (which may span multiple lines) and return
    a flat list of ``(line_text, is_bold)`` tuples.

    The split is applied to the ENTIRE block first so that::

        **Chorus line 1
        Chorus line 2**

    produces TWO bold lines, not just the first.

    *base_bold=True* makes every segment bold by default (used for
    the separate chorus_text block).
    """
    parts = re.split(r'\*\*', text, flags=re.DOTALL)
    result: list[tuple[str, bool]] = []
    for i, part in enumerate(parts):
        if not part:          # skip empty strings (leading/trailing **)
            continue
        block_bold = base_bold or (i % 2 == 1)
        for line in part.split('\n'):
            result.append((line, block_bold))
    return result


def _make_rPr(saved_rPr, is_bold: bool, font_size_pt: int | None = None):
    """
    Return a fresh ``<a:rPr>`` lxml element.

    Starts from a deep-copy of *saved_rPr* (the template's run properties)
    so font face and colour are preserved.  Bold and optional font size are
    applied on top.

    *font_size_pt* is in points; OOXML stores it in 1/100 pt (``sz`` attr).
    When *None*, the template's original size is kept.
    """
    if saved_rPr is not None:
        rPr = copy.deepcopy(saved_rPr)
    else:
        rPr = etree.Element(f"{{{_NS}}}rPr")
        rPr.set("lang", "en-US")
        rPr.set("dirty", "0")
    rPr.set("b", "1" if is_bold else "0")
    if font_size_pt is not None:
        rPr.set("sz", str(int(font_size_pt * 100)))   # e.g. 36pt → "3600"
    return rPr


def _set_frame_text(shape, verse_text: str, chorus_text: str = ""):
    """
    Replace the entire content of *shape*'s text frame with the supplied
    verse and optional chorus, with:

    * **Multi-line bold**: ``**...**`` delimiters are scanned across the
      WHOLE text block first (not line-by-line), so a chorus that spans
      several lines is correctly bolded throughout.
    * **Smart font scaling**: font size steps down based on total character
      count so long lyrics never overflow the text box.
    * **Margin compression**: inner padding is tightened for long lyrics to
      give the text box extra room.
    * **Template fidelity**: font face and colour are preserved from the
      template's first run.
    """
    text_frame = shape.text_frame

    # 1. Normalise escaped \\n sequences
    verse_text  = _normalise(verse_text)
    chorus_text = _normalise(chorus_text) if chorus_text else ""

    # 2. Compute total printable characters for font-size decision
    combined    = verse_text + ("\n" + chorus_text if chorus_text else "")
    total_chars = len(combined.replace("\n", ""))

    # 3. Smart step-down font size (thresholds tuned by user testing)
    #    < 258 chars  → 36 pt  (comfortable)
    #    258–367      → 34 pt  (slightly denser)
    #    368–369      → 32 pt  (boundary zone)
    #    ≥ 370        → 28 pt  (dense slide)
    if total_chars >= 370:
        font_size = 28
    elif total_chars >= 368:
        font_size = 32
    elif total_chars >= 258:
        font_size = 34
    else:
        font_size = 36

    # 4. Word wrap — required for all layouts
    text_frame.word_wrap = True

    # 5. Compress margins for long lyrics so text has more room
    if total_chars > 300:
        text_frame.margin_left   = Emu(45720)   # 0.05"
        text_frame.margin_right  = Emu(45720)
        text_frame.margin_top    = Emu(27432)   # 0.03"
        text_frame.margin_bottom = Emu(27432)

    # 6. Snapshot template run properties before clearing
    saved_rPr = _snapshot_rPr(text_frame)

    # 7. Build flat list of (line_text, is_bold) by parsing ** blocks
    #    across the entire text first (handles multi-line bold blocks)
    lines: list[tuple[str, bool]] = _parse_lines_with_bold(verse_text, base_bold=False)
    if chorus_text:
        lines.append(("", False))          # blank separator
        lines.extend(_parse_lines_with_bold(chorus_text, base_bold=True))

    # 8. Clear existing <a:p> elements
    txBody = text_frame._txBody
    for p_el in txBody.findall(f"{{{_NS}}}p"):
        txBody.remove(p_el)

    # 9. Rebuild — one <a:p> per line
    #    Each (non-empty) line may contain multiple <a:r> runs if it has
    #    inline ** markers that weren't already consumed by block-level parsing.
    for line_text, line_is_bold in lines:
        p_el = etree.SubElement(txBody, f"{{{_NS}}}p")

        # Inline ** pass (handles cases like "Verse **word** more verse")
        segments = _parse_lines_with_bold(line_text, base_bold=line_is_bold)

        if not segments:
            # Completely empty line — single empty run preserves the blank row
            r_el = etree.SubElement(p_el, f"{{{_NS}}}r")
            r_el.insert(0, _make_rPr(saved_rPr, line_is_bold, font_size))
            t_el = etree.SubElement(r_el, f"{{{_NS}}}t")
            t_el.text = ""
        else:
            for seg_text, seg_is_bold in segments:
                r_el = etree.SubElement(p_el, f"{{{_NS}}}r")
                r_el.insert(0, _make_rPr(saved_rPr, seg_is_bold, font_size))
                t_el = etree.SubElement(r_el, f"{{{_NS}}}t")
                t_el.text = seg_text


# ---------------------------------------------------------------------------
# Background image copying
# ---------------------------------------------------------------------------

# PresentationML namespace
_PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"
# Relationship type for embedded images
_IMG_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
# r: namespace used in blip embed attributes
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _copy_background(src_slide, dst_slide):
    """
    Copy the explicit slide background (<p:bg>) from *src_slide* to
    *dst_slide*, including re-registering any embedded image relationships
    so the background image renders correctly on the new slide.

    In OOXML, a slide's background lives inside <p:cSld><p:bg> — a sibling
    of <p:spTree> — so copying the shape tree alone silently drops it.
    """
    # Locate <p:bg> inside the source slide's <p:cSld>
    src_bg = src_slide._element.find(f"{{{_PNS}}}cSld/{{{_PNS}}}bg")
    if src_bg is None:
        return  # This slide has no explicit background; layout bg is inherited

    # Deep-copy the <p:bg> XML node (safe: lxml elements have no back-refs)
    new_bg = copy.deepcopy(src_bg)

    # Re-map every <a:blip r:embed="rIdX"/> reference.
    # Each rId points to an image Part registered on the SOURCE slide's part.
    # We must add the same image Part to the DESTINATION slide's part and
    # update the rId in the copied XML to the new destination rId.
    for blip in new_bg.findall(f".//{{{_NS}}}blip"):
        old_rId = blip.get(f"{{{_R_NS}}}embed")
        if old_rId:
            try:
                img_part = src_slide.part.related_part(old_rId)
                new_rId  = dst_slide.part.relate_to(img_part, _IMG_REL)
                blip.set(f"{{{_R_NS}}}embed", new_rId)
            except Exception as exc:
                print(f"[generate-ppt] Warning: could not copy bg image rId {old_rId}: {exc}",
                      flush=True)

    # Insert / replace <p:bg> in the destination slide's <p:cSld>
    dst_cSld = dst_slide._element.find(f"{{{_PNS}}}cSld")
    if dst_cSld is not None:
        existing = dst_cSld.find(f"{{{_PNS}}}bg")
        if existing is not None:
            dst_cSld.remove(existing)
        # <p:bg> must be the FIRST child of <p:cSld> per the OOXML schema
        dst_cSld.insert(0, new_bg)


# ---------------------------------------------------------------------------
# Safe slide duplication
# ---------------------------------------------------------------------------

def _duplicate_slide(prs: Presentation, slide_index: int):
    """
    Safely duplicate the slide at *slide_index* and insert the copy
    immediately after it in the deck.

    Algorithm:
      1. Identify the slide layout used by the template slide.
      2. Add a fresh slide at the end using that same layout.
      3. Copy the source slide's spTree (shapes) into the new slide.
      4. Copy the source slide's explicit background (<p:bg>) including
         re-registering any embedded image relationships.
      5. Reposition the new slide immediately after the source.

    Returns the new slide object (now at position slide_index + 1).
    """
    src_slide = prs.slides[slide_index]
    layout    = src_slide.slide_layout

    # Step 1: Create a properly initialised slide (appended at the end)
    new_slide = prs.slides.add_slide(layout)

    # Step 2: Copy shape tree
    src_spTree = src_slide.shapes._spTree
    dst_spTree = new_slide.shapes._spTree

    for child in list(dst_spTree):
        dst_spTree.remove(child)
    for child in src_spTree:
        dst_spTree.append(copy.deepcopy(child))

    # Step 3: Copy background image (the key fix — spTree copy skips <p:bg>)
    _copy_background(src_slide, new_slide)

    # Step 4: Reposition the new slide immediately after the source
    sldIdLst  = prs.slides._sldIdLst
    new_entry = sldIdLst[-1]
    sldIdLst.remove(new_entry)
    sldIdLst.insert(slide_index + 1, new_entry)

    return prs.slides[slide_index + 1]


# ---------------------------------------------------------------------------
# Section slide filling
# ---------------------------------------------------------------------------

def _fill_section_slide(prs: Presentation, slide_index: int, song: dict):
    """
    Fill a Mass-section slide with the song's lyrics, duplicating the slide
    for each additional verse.

    Key fix: the lyrics shape is identified ONCE by its {{LYRICS}} placeholder
    *before* any content is written. Subsequent duplicated slides use the same
    shape index so we never need to search for the placeholder again (it will
    have been replaced by verse 1 text by the time we duplicate).
    """
    lyrics      = song.get("lyrics", [])
    chorus_text = ""
    verses: list[str] = []

    for item in lyrics:
        label = item.get("label", "").lower()
        text  = item.get("text", "").strip()
        if "chorus" in label:
            chorus_text = text
        elif text:
            verses.append(text)

    # Edge case: only a chorus provided — treat it as the sole verse
    if not verses and chorus_text:
        verses      = [chorus_text]
        chorus_text = ""

    # Find the lyrics shape index BEFORE modifying anything
    original_slide    = prs.slides[slide_index]
    lyrics_shape_idx  = None
    shapes_list       = list(original_slide.shapes)

    for idx, shape in enumerate(shapes_list):
        if shape.has_text_frame and _frame_contains(shape.text_frame, "{{LYRICS}}"):
            lyrics_shape_idx = idx
            break

    # If no lyrics placeholder found at all, clear it and bail
    if lyrics_shape_idx is None:
        for shape in shapes_list:
            if shape.has_text_frame and _frame_contains(shape.text_frame, "{{LYRICS}}"):
                _replace_text_in_frame(shape.text_frame, "{{LYRICS}}", "")
        return

    if not verses:
        # No lyrics — clear placeholder
        _replace_text_in_frame(
            shapes_list[lyrics_shape_idx].text_frame, "{{LYRICS}}", ""
        )
        return

    # --- Fill Verse 1 on the original slide ---
    _set_frame_text(
        shapes_list[lyrics_shape_idx],
        verses[0],
        chorus_text,
    )

    # --- Duplicate slide for each additional verse ---
    for extra_idx, verse_text in enumerate(verses[1:], start=1):
        # The slide we want to duplicate is now at (slide_index + extra_idx - 1)
        _duplicate_slide(prs, slide_index + extra_idx - 1)

        # The new duplicate is immediately after: slide_index + extra_idx
        dup_slide    = prs.slides[slide_index + extra_idx]
        dup_shapes   = list(dup_slide.shapes)

        # Use the same shape index — no need to search for {{LYRICS}}
        if lyrics_shape_idx < len(dup_shapes):
            _set_frame_text(
                dup_shapes[lyrics_shape_idx],
                verse_text,
                chorus_text,
            )


# ---------------------------------------------------------------------------
# Core generation
# ---------------------------------------------------------------------------

def generate_presentation(template_url: str, date: str, sections: list) -> io.BytesIO:
    """
    Download the PPTX template, fill {{DATE}} and {{LYRICS}} placeholders,
    duplicate slides for multi-verse hymns, and return a BytesIO buffer.
    """
    # 1. Download template entirely into memory (no temp files)
    resp = requests.get(template_url, timeout=30)
    resp.raise_for_status()
    template_bytes = io.BytesIO(resp.content)

    # 2. Open presentation
    prs = Presentation(template_bytes)

    # 3. Cover page (slide 0) — replace {{DATE}}
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            _replace_text_in_frame(shape.text_frame, "{{DATE}}", _normalise(date))

    # 4. Build section name → song dict lookup (case-insensitive)
    section_map = {
        s.get("name", "").strip().lower(): s.get("song", {})
        for s in sections
    }

    # 5. Fill each Mass-section slide, tracking index offset from duplications
    slide_offset = 0
    for section_idx, section_name in enumerate(SECTION_ORDER):
        current_index = 1 + section_idx + slide_offset
        song          = section_map.get(section_name.lower(), {})

        before = len(prs.slides)
        _fill_section_slide(prs, current_index, song)
        slide_offset += len(prs.slides) - before

    # 6. Save to in-memory buffer
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ---------------------------------------------------------------------------
# Chord PDF merging (server-side — bypasses all browser CORS/binary issues)
# ---------------------------------------------------------------------------

_DRIVE_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/124.0 Safari/537.36"
    ),
    "Accept": "application/pdf,*/*;q=0.8",
}


def _drive_download_url(view_url: str) -> str:
    """
    Convert any Google Drive sharing / view URL to a direct download URL.
    Supports:
      https://drive.google.com/file/d/FILE_ID/view
      https://drive.google.com/open?id=FILE_ID
      https://drive.google.com/uc?id=FILE_ID  (already a download-ish URL)
    """
    import re
    m = re.search(r'/d/([a-zA-Z0-9_-]+)', view_url)
    if not m:
        m = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', view_url)
    if m:
        return f"https://drive.google.com/uc?export=download&id={m.group(1)}"
    return view_url  # Not a Drive URL — pass through


def _fetch_pdf_bytes(session: requests.Session, url: str) -> bytes:
    """
    Fetch PDF bytes from *url* using *session*, handling Google Drive's
    virus-scan confirmation page automatically.
    Raises ValueError if the final response is not a valid PDF.
    """
    resp = session.get(url, headers=_DRIVE_HEADERS, timeout=30, allow_redirects=True)
    resp.raise_for_status()

    # If Google returns HTML (virus-scan warning), extract confirm token and retry
    if resp.content[:4] != b"%PDF":
        confirm_token = "t"  # works for most small public files
        for key, value in resp.cookies.items():
            if "download_warning" in key.lower():
                confirm_token = value
                break
        sep = "&" if "?" in url else "?"
        confirmed = f"{url}{sep}confirm={confirm_token}"
        print(f"[merge-chords] Retrying with confirm={confirm_token}", flush=True)
        resp = session.get(confirmed, headers=_DRIVE_HEADERS, timeout=30, allow_redirects=True)
        resp.raise_for_status()

    if resp.content[:4] != b"%PDF":
        snippet = resp.content[:120].decode("utf-8", errors="replace")
        raise ValueError(
            f"Response is not a PDF (starts with: {snippet!r}). "
            "Make sure the file is shared as 'Anyone with the link can view'."
        )

    return resp.content


def merge_chord_pdfs(sections: list) -> io.BytesIO:
    """
    Download every chord PDF in *sections* order and merge them into one PDF.

    *sections* is a list of {"name": str, "url": str} dicts.
    Returns an in-memory BytesIO buffer of the merged PDF.
    Raises RuntimeError if no valid PDFs could be loaded.
    """
    writer = PdfWriter()
    session = requests.Session()
    errors = []

    for entry in sections:
        name = entry.get("name", "?")
        raw_url = entry.get("url", "").strip()
        if not raw_url:
            continue
        dl_url = _drive_download_url(raw_url)
        print(f"[merge-chords] Fetching '{name}': {dl_url[:80]}", flush=True)
        try:
            pdf_bytes = _fetch_pdf_bytes(session, dl_url)
            reader = PdfReader(io.BytesIO(pdf_bytes))
            for page in reader.pages:
                writer.add_page(page)
            print(f"[merge-chords] '{name}' — {len(reader.pages)} page(s) added", flush=True)
        except Exception as exc:
            errors.append(f"{name}: {exc}")
            print(f"[merge-chords] SKIP '{name}': {exc}", flush=True)

    if writer.get_num_pages() == 0:
        raise RuntimeError(
            "No chord PDFs could be loaded.\n" + "\n".join(errors)
        )

    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# HTTP handler — Vercel-compatible + local http.server
# ---------------------------------------------------------------------------

class handler(BaseHTTPRequestHandler):

    def log_message(self, fmt, *args):  # suppress default logging
        print(f"[{self.client_address[0]}] {fmt % args}", flush=True)

    # ------------------------------------------------------------------
    def do_OPTIONS(self):  # noqa: N802  (CORS preflight)
        self.send_response(204)
        self._cors_headers()
        self.end_headers()

    def do_GET(self):  # noqa: N802
        parsed = urlparse(self.path)

        if parsed.path in ("/proxy-pdf", "/api/proxy-pdf"):
            self._handle_proxy_pdf(parse_qs(parsed.query))
        else:
            self._send_text(200, "generate-ppt is running. POST JSON to /generate-ppt.")

    def _handle_proxy_pdf(self, params: dict):
        """
        Server-side PDF proxy: fetches a remote PDF URL and returns the bytes.
        Sidesteps browser CORS restrictions; all Google Drive traffic is
        server-to-server.

        Handles Google Drive's virus-scan confirmation page by:
          1. Making an initial request with a Session (captures cookies).
          2. If the response is HTML and a 'download_warning' cookie exists,
             re-requesting with confirm=t to bypass the interstitial.
          3. Validating the %%PDF magic bytes before sending to the client.
        """
        raw_url = params.get("url", [""])[0].strip()
        if not raw_url:
            self._send_json_error(400, "Missing 'url' query parameter")
            return

        target_url = unquote(raw_url)
        print(f"[proxy-pdf] Fetching: {target_url[:100]}", flush=True)

        _HEADERS = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0 Safari/537.36"
            ),
            "Accept": "application/pdf,*/*;q=0.8",
        }

        try:
            session = requests.Session()

            # --- First attempt ---
            resp = session.get(
                target_url, headers=_HEADERS, timeout=30, allow_redirects=True
            )
            resp.raise_for_status()

            content_type = resp.headers.get("Content-Type", "")

            # --- Handle Google Drive confirmation page ---
            # For files > ~25 MB Google returns an HTML page with a warning cookie.
            # We detect this either via Content-Type or by checking the magic bytes.
            is_html = "text/html" in content_type
            is_pdf_bytes = resp.content[:4] == b"%PDF"

            if is_html or not is_pdf_bytes:
                # Try to extract the confirmation token from cookies
                confirm_token = None
                for key, value in resp.cookies.items():
                    if "download_warning" in key.lower():
                        confirm_token = value
                        break

                if not confirm_token:
                    # Fallback: try confirm=t which works for most public files
                    confirm_token = "t"

                sep = "&" if "?" in target_url else "?"
                confirmed_url = f"{target_url}{sep}confirm={confirm_token}"
                print(f"[proxy-pdf] Retrying with confirm token: {confirm_token}", flush=True)

                resp = session.get(
                    confirmed_url, headers=_HEADERS, timeout=30, allow_redirects=True
                )
                resp.raise_for_status()

            # --- Final magic-byte validation ---
            if resp.content[:4] != b"%PDF":
                snippet = resp.content[:200].decode("utf-8", errors="replace")
                print(f"[proxy-pdf] Non-PDF response snippet: {snippet}", flush=True)
                self._send_json_error(
                    502,
                    "Google Drive did not return a PDF. "
                    "Ensure the file is shared as 'Anyone with the link can view' "
                    "and the chord_link column contains a valid Google Drive URL.",
                )
                return

            data = resp.content
            print(f"[proxy-pdf] OK — {len(data):,} bytes", flush=True)

            self.send_response(200)
            self.send_header("Content-Type", "application/pdf")
            self.send_header("Content-Length", str(len(data)))
            self.send_header("Content-Disposition", "inline")
            self._cors_headers()
            self.end_headers()
            self.wfile.write(data)

        except requests.HTTPError as exc:
            msg = f"Remote server returned an error fetching the PDF: {exc}"
            print(f"[proxy-pdf] ERROR: {msg}", flush=True)
            self._send_json_error(502, msg)
        except Exception:
            tb = traceback.format_exc()
            print(f"[proxy-pdf] TRACEBACK:\n{tb}", flush=True)
            self._send_json_error(500, f"Proxy failed: {tb}")


    def do_POST(self):  # noqa: N802
        length = int(self.headers.get("Content-Length", 0))
        raw    = self.rfile.read(length)

        try:
            body = json.loads(raw)
        except json.JSONDecodeError as exc:
            self._send_json_error(400, f"Invalid JSON: {exc}")
            return

        # ---- Route by body content (immune to proxy path rewriting) --------
        # merge-chords requests send { "sections": [{name, url}] } — no template_url
        # generate-ppt requests always include template_url
        if "template_url" not in body:
            self._handle_merge_chords(body)
            return

        # ---- Default: generate PPT ----------------------------------------
        template_url = body.get("template_url", "").strip()
        date         = body.get("date", "")
        sections     = body.get("sections", [])

        if not template_url:
            self._send_json_error(400, "Missing required field: template_url")
            return

        print(f"[generate-ppt] date='{date}'  sections={len(sections)}", flush=True)

        try:
            buf = generate_presentation(template_url, date, sections)
        except requests.HTTPError as exc:
            msg = f"Failed to download template: {exc}"
            print(f"[generate-ppt] ERROR: {msg}", flush=True)
            self._send_json_error(502, msg)
            return
        except Exception:
            tb = traceback.format_exc()
            print(f"[generate-ppt] TRACEBACK:\n{tb}", flush=True)
            self._send_json_error(500, f"Presentation generation failed:\n{tb}")
            return

        data = buf.read()
        print(f"[generate-ppt] OK — {len(data):,} bytes", flush=True)

        self.send_response(200)
        self.send_header(
            "Content-Type",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.send_header("Content-Disposition", 'attachment; filename="Mass_Presentation.pptx"')
        self.send_header("Content-Length", str(len(data)))
        self._cors_headers()
        self.end_headers()
        self.wfile.write(data)

    # ------------------------------------------------------------------
    def _handle_merge_chords(self, body: dict):
        """
        POST /merge-chords handler.
        Accepts { "sections": [{"name": str, "url": str}, ...] },
        downloads and merges all chord PDFs server-side using pypdf,
        and returns the combined PDF binary.
        """
        sections = body.get("sections", [])
        if not sections:
            self._send_json_error(400, "No sections provided.")
            return

        print(f"[merge-chords] Merging {len(sections)} chord sheet(s)", flush=True)
        try:
            buf  = merge_chord_pdfs(sections)
            data = buf.read()
        except RuntimeError as exc:
            self._send_json_error(502, str(exc))
            return
        except Exception:
            tb = traceback.format_exc()
            print(f"[merge-chords] TRACEBACK:\n{tb}", flush=True)
            self._send_json_error(500, f"Chord merge failed:\n{tb}")
            return

        print(f"[merge-chords] Done — {len(data):,} bytes", flush=True)
        self.send_response(200)
        self.send_header("Content-Type", "application/pdf")
        self.send_header("Content-Disposition", 'attachment; filename="Mass_Chords.pdf"')
        self.send_header("Content-Length", str(len(data)))
        self._cors_headers()
        self.end_headers()
        self.wfile.write(data)

    # ------------------------------------------------------------------
    def _cors_headers(self):
        self.send_header("Access-Control-Allow-Origin",  "*")
        self.send_header("Access-Control-Allow-Methods", "POST, GET, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def _send_json_error(self, code: int, message: str):
        body = json.dumps({"error": message}).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self._cors_headers()
        self.end_headers()
        self.wfile.write(body)

    def _send_text(self, code: int, message: str):
        body = message.encode()
        self.send_response(code)
        self.send_header("Content-Type", "text/plain")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)


# ---------------------------------------------------------------------------
# Local dev entry-point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    from http.server import HTTPServer

    port = 8000
    print(f"[generate-ppt] Starting -> http://localhost:{port}", flush=True)
    server = HTTPServer(("", port), handler)
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n[generate-ppt] Shutting down.", flush=True)
