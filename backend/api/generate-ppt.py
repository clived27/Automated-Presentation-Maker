"""
generate-ppt.py
Serverless function: accepts POST with JSON body, downloads a master PPTX
template, fills placeholders, and returns a binary .pptx download.

Expected JSON body:
{
    "template_url": "https://...",
    "date": "Sunday, 29 June 2026",
    "sections": [
        {
            "name": "Entrance",
            "song": {
                "title": "Song Title",
                "lyrics": [
                    { "label": "Verse 1", "text": "Line 1\\nLine 2" },
                    { "label": "Chorus",  "text": "Chorus line 1\\nChorus line 2" },
                    { "label": "Verse 2", "text": "Line 3\\nLine 4" }
                ]
            }
        },
        ... (9 sections total)
    ]
}

Slide mapping (0-indexed):
  0  = Cover Page     (replace {{DATE}})
  1  = Entrance
  2  = Lord Have Mercy
  3  = Gloria
  4  = Acclamation
  5  = Offertory
  6  = Holy Holy
  7  = Proclamation
  8  = Communion
  9  = Recessional
"""

import copy
import json
import io
import traceback
from http.server import BaseHTTPRequestHandler

import requests
from pptx import Presentation
from pptx.util import Pt
from lxml import etree


# ---------------------------------------------------------------------------
# OOXML DrawingML namespace
# ---------------------------------------------------------------------------
_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Section order (must match slide indices 1-9 in the template)
# ---------------------------------------------------------------------------
SECTION_ORDER = [
    "Entrance",
    "Lord Have Mercy",
    "Gloria",
    "Acclamation",
    "Offertory",
    "Holy Holy",
    "Proclamation",
    "Communion",
    "Recessional",
]


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def _normalise(text: str) -> str:
    """
    Convert literal backslash-n sequences (as stored in some databases)
    into real newline characters, and strip surrounding whitespace.
    """
    return text.replace("\\n", "\n").strip()


def _para_full_text(para) -> str:
    """Concatenate all run texts in a paragraph (handles split-run placeholders)."""
    return "".join(run.text for run in para.runs)


def _frame_contains(text_frame, placeholder: str) -> bool:
    """
    Return True if *placeholder* exists anywhere in the text frame.
    Checks concatenated paragraph text so split-run placeholders are found.
    """
    for para in text_frame.paragraphs:
        if placeholder in _para_full_text(para):
            return True
    return False


def _replace_in_para(para, old: str, new: str) -> bool:
    """
    Replace *old* with *new* within a paragraph, even if the placeholder
    is split across multiple runs (consolidates into the first run).
    """
    full = _para_full_text(para)
    if old not in full:
        return False
    replaced = full.replace(old, new)
    runs = para.runs
    if runs:
        runs[0].text = replaced
        for run in runs[1:]:
            run.text = ""
    return True


def _replace_text_in_frame(text_frame, old: str, new: str) -> bool:
    """Replace *old* with *new* across every paragraph in a text frame."""
    changed = False
    for para in text_frame.paragraphs:
        if _replace_in_para(para, old, new):
            changed = True
    return changed


def _snapshot_rPr(text_frame):
    """
    Return a deep-copied <a:rPr> element from the first run in the frame,
    or None if the frame has no runs with run properties.
    Used to carry over font name, size, colour when building new paragraphs.
    """
    for para in text_frame.paragraphs:
        for run in para.runs:
            rPr_el = run._r.find(f"{{{_NS}}}rPr")
            if rPr_el is not None:
                return copy.deepcopy(rPr_el)
    return None


def _set_frame_text(text_frame, verse_text: str, chorus_text: str = ""):
    """
    Replace the entire content of *text_frame* with the supplied verse and
    optional chorus.

    Features:
    - Normalises escaped \\n sequences into real newlines.
    - Enables word-wrap on the text frame.
    - Renders chorus lines in BOLD, separated from the verse by a blank line.
    - Preserves the template's original run properties (font face, colour).
    """
    # 1. Normalise text
    verse_text  = _normalise(verse_text)
    chorus_text = _normalise(chorus_text) if chorus_text else ""

    # 2. Enable word wrap so long lines don't escape the text box boundary
    text_frame.word_wrap = True

    # 3. Snapshot the template run properties before clearing anything
    saved_rPr = _snapshot_rPr(text_frame)

    # 4. Build a list of (line_text, is_bold) tuples
    lines: list[tuple[str, bool]] = []
    for line in verse_text.split("\n"):
        lines.append((line, False))
    if chorus_text:
        lines.append(("", False))       # blank separator before chorus
        for line in chorus_text.split("\n"):
            lines.append((line, True))  # chorus lines → bold

    # 5. Clear existing <a:p> elements from the txBody
    txBody = text_frame._txBody
    for p_el in txBody.findall(f"{{{_NS}}}p"):
        txBody.remove(p_el)

    # 6. Rebuild one <a:p> per line
    for text, is_bold in lines:
        p_el = etree.SubElement(txBody, f"{{{_NS}}}p")
        r_el = etree.SubElement(p_el, f"{{{_NS}}}r")

        # Build <a:rPr> — start from snapshot or create fresh
        if saved_rPr is not None:
            rPr = copy.deepcopy(saved_rPr)
        else:
            rPr = etree.Element(f"{{{_NS}}}rPr")
            rPr.set("lang", "en-US")
            rPr.set("dirty", "0")

        # Apply bold for chorus lines (template font size preserved)
        rPr.set("b", "1" if is_bold else "0")

        r_el.insert(0, rPr)

        t_el = etree.SubElement(r_el, f"{{{_NS}}}t")
        t_el.text = text


# ---------------------------------------------------------------------------
# Background image copying
# ---------------------------------------------------------------------------

# PresentationML namespace
_PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"
# Relationship type for embedded images
_IMG_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
# r: namespace used in blip embed attributes
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _copy_background(src_slide, dst_slide):
    """
    Copy the explicit slide background (<p:bg>) from *src_slide* to
    *dst_slide*, including re-registering any embedded image relationships
    so the background image renders correctly on the new slide.

    In OOXML, a slide's background lives inside <p:cSld><p:bg> — a sibling
    of <p:spTree> — so copying the shape tree alone silently drops it.
    """
    # Locate <p:bg> inside the source slide's <p:cSld>
    src_bg = src_slide._element.find(f"{{{_PNS}}}cSld/{{{_PNS}}}bg")
    if src_bg is None:
        return  # This slide has no explicit background; layout bg is inherited

    # Deep-copy the <p:bg> XML node (safe: lxml elements have no back-refs)
    new_bg = copy.deepcopy(src_bg)

    # Re-map every <a:blip r:embed="rIdX"/> reference.
    # Each rId points to an image Part registered on the SOURCE slide's part.
    # We must add the same image Part to the DESTINATION slide's part and
    # update the rId in the copied XML to the new destination rId.
    for blip in new_bg.findall(f".//{{{_NS}}}blip"):
        old_rId = blip.get(f"{{{_R_NS}}}embed")
        if old_rId:
            try:
                img_part = src_slide.part.related_part(old_rId)
                new_rId  = dst_slide.part.relate_to(img_part, _IMG_REL)
                blip.set(f"{{{_R_NS}}}embed", new_rId)
            except Exception as exc:
                print(f"[generate-ppt] Warning: could not copy bg image rId {old_rId}: {exc}",
                      flush=True)

    # Insert / replace <p:bg> in the destination slide's <p:cSld>
    dst_cSld = dst_slide._element.find(f"{{{_PNS}}}cSld")
    if dst_cSld is not None:
        existing = dst_cSld.find(f"{{{_PNS}}}bg")
        if existing is not None:
            dst_cSld.remove(existing)
        # <p:bg> must be the FIRST child of <p:cSld> per the OOXML schema
        dst_cSld.insert(0, new_bg)


# ---------------------------------------------------------------------------
# Safe slide duplication
# ---------------------------------------------------------------------------

def _duplicate_slide(prs: Presentation, slide_index: int):
    """
    Safely duplicate the slide at *slide_index* and insert the copy
    immediately after it in the deck.

    Algorithm:
      1. Identify the slide layout used by the template slide.
      2. Add a fresh slide at the end using that same layout.
      3. Copy the source slide's spTree (shapes) into the new slide.
      4. Copy the source slide's explicit background (<p:bg>) including
         re-registering any embedded image relationships.
      5. Reposition the new slide immediately after the source.

    Returns the new slide object (now at position slide_index + 1).
    """
    src_slide = prs.slides[slide_index]
    layout    = src_slide.slide_layout

    # Step 1: Create a properly initialised slide (appended at the end)
    new_slide = prs.slides.add_slide(layout)

    # Step 2: Copy shape tree
    src_spTree = src_slide.shapes._spTree
    dst_spTree = new_slide.shapes._spTree

    for child in list(dst_spTree):
        dst_spTree.remove(child)
    for child in src_spTree:
        dst_spTree.append(copy.deepcopy(child))

    # Step 3: Copy background image (the key fix — spTree copy skips <p:bg>)
    _copy_background(src_slide, new_slide)

    # Step 4: Reposition the new slide immediately after the source
    sldIdLst  = prs.slides._sldIdLst
    new_entry = sldIdLst[-1]
    sldIdLst.remove(new_entry)
    sldIdLst.insert(slide_index + 1, new_entry)

    return prs.slides[slide_index + 1]


# ---------------------------------------------------------------------------
# Section slide filling
# ---------------------------------------------------------------------------

def _fill_section_slide(prs: Presentation, slide_index: int, song: dict):
    """
    Fill a Mass-section slide with the song's lyrics, duplicating the slide
    for each additional verse.

    Key fix: the lyrics shape is identified ONCE by its {{LYRICS}} placeholder
    *before* any content is written. Subsequent duplicated slides use the same
    shape index so we never need to search for the placeholder again (it will
    have been replaced by verse 1 text by the time we duplicate).
    """
    lyrics      = song.get("lyrics", [])
    chorus_text = ""
    verses: list[str] = []

    for item in lyrics:
        label = item.get("label", "").lower()
        text  = item.get("text", "").strip()
        if "chorus" in label:
            chorus_text = text
        elif text:
            verses.append(text)

    # Edge case: only a chorus provided — treat it as the sole verse
    if not verses and chorus_text:
        verses      = [chorus_text]
        chorus_text = ""

    # Find the lyrics shape index BEFORE modifying anything
    original_slide    = prs.slides[slide_index]
    lyrics_shape_idx  = None
    shapes_list       = list(original_slide.shapes)

    for idx, shape in enumerate(shapes_list):
        if shape.has_text_frame and _frame_contains(shape.text_frame, "{{LYRICS}}"):
            lyrics_shape_idx = idx
            break

    # If no lyrics placeholder found at all, clear it and bail
    if lyrics_shape_idx is None:
        for shape in shapes_list:
            if shape.has_text_frame and _frame_contains(shape.text_frame, "{{LYRICS}}"):
                _replace_text_in_frame(shape.text_frame, "{{LYRICS}}", "")
        return

    if not verses:
        # No lyrics — clear placeholder
        _replace_text_in_frame(
            shapes_list[lyrics_shape_idx].text_frame, "{{LYRICS}}", ""
        )
        return

    # --- Fill Verse 1 on the original slide ---
    _set_frame_text(
        shapes_list[lyrics_shape_idx].text_frame,
        verses[0],
        chorus_text,
    )

    # --- Duplicate slide for each additional verse ---
    for extra_idx, verse_text in enumerate(verses[1:], start=1):
        # The slide we want to duplicate is now at (slide_index + extra_idx - 1)
        _duplicate_slide(prs, slide_index + extra_idx - 1)

        # The new duplicate is immediately after: slide_index + extra_idx
        dup_slide    = prs.slides[slide_index + extra_idx]
        dup_shapes   = list(dup_slide.shapes)

        # Use the same shape index — no need to search for {{LYRICS}}
        if lyrics_shape_idx < len(dup_shapes):
            _set_frame_text(
                dup_shapes[lyrics_shape_idx].text_frame,
                verse_text,
                chorus_text,
            )


# ---------------------------------------------------------------------------
# Core generation
# ---------------------------------------------------------------------------

def generate_presentation(template_url: str, date: str, sections: list) -> io.BytesIO:
    """
    Download the PPTX template, fill {{DATE}} and {{LYRICS}} placeholders,
    duplicate slides for multi-verse hymns, and return a BytesIO buffer.
    """
    # 1. Download template entirely into memory (no temp files)
    resp = requests.get(template_url, timeout=30)
    resp.raise_for_status()
    template_bytes = io.BytesIO(resp.content)

    # 2. Open presentation
    prs = Presentation(template_bytes)

    # 3. Cover page (slide 0) — replace {{DATE}}
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            _replace_text_in_frame(shape.text_frame, "{{DATE}}", _normalise(date))

    # 4. Build section name → song dict lookup (case-insensitive)
    section_map = {
        s.get("name", "").strip().lower(): s.get("song", {})
        for s in sections
    }

    # 5. Fill each Mass-section slide, tracking index offset from duplications
    slide_offset = 0
    for section_idx, section_name in enumerate(SECTION_ORDER):
        current_index = 1 + section_idx + slide_offset
        song          = section_map.get(section_name.lower(), {})

        before = len(prs.slides)
        _fill_section_slide(prs, current_index, song)
        slide_offset += len(prs.slides) - before

    # 6. Save to in-memory buffer
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ---------------------------------------------------------------------------
# HTTP handler — Vercel-compatible + local http.server
# ---------------------------------------------------------------------------

class handler(BaseHTTPRequestHandler):

    def log_message(self, fmt, *args):  # suppress default logging
        print(f"[{self.client_address[0]}] {fmt % args}", flush=True)

    # ------------------------------------------------------------------
    def do_OPTIONS(self):  # noqa: N802  (CORS preflight)
        self.send_response(204)
        self._cors_headers()
        self.end_headers()

    def do_GET(self):  # noqa: N802
        self._send_text(200, "generate-ppt is running. POST JSON to this endpoint.")

    def do_POST(self):  # noqa: N802
        length  = int(self.headers.get("Content-Length", 0))
        raw     = self.rfile.read(length)

        try:
            body = json.loads(raw)
        except json.JSONDecodeError as exc:
            self._send_json_error(400, f"Invalid JSON: {exc}")
            return

        template_url = body.get("template_url", "").strip()
        date         = body.get("date", "")
        sections     = body.get("sections", [])

        if not template_url:
            self._send_json_error(400, "Missing required field: template_url")
            return

        print(f"[generate-ppt] date='{date}'  sections={len(sections)}", flush=True)

        try:
            buf = generate_presentation(template_url, date, sections)
        except requests.HTTPError as exc:
            msg = f"Failed to download template: {exc}"
            print(f"[generate-ppt] ERROR: {msg}", flush=True)
            self._send_json_error(502, msg)
            return
        except Exception:
            tb = traceback.format_exc()
            print(f"[generate-ppt] TRACEBACK:\n{tb}", flush=True)
            self._send_json_error(500, f"Presentation generation failed:\n{tb}")
            return

        data = buf.read()
        print(f"[generate-ppt] OK — {len(data):,} bytes", flush=True)

        self.send_response(200)
        self.send_header(
            "Content-Type",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.send_header("Content-Disposition", 'attachment; filename="Mass_Presentation.pptx"')
        self.send_header("Content-Length", str(len(data)))
        self._cors_headers()
        self.end_headers()
        self.wfile.write(data)

    # ------------------------------------------------------------------
    def _cors_headers(self):
        self.send_header("Access-Control-Allow-Origin",  "*")
        self.send_header("Access-Control-Allow-Methods", "POST, GET, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def _send_json_error(self, code: int, message: str):
        body = json.dumps({"error": message}).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self._cors_headers()
        self.end_headers()
        self.wfile.write(body)

    def _send_text(self, code: int, message: str):
        body = message.encode()
        self.send_response(code)
        self.send_header("Content-Type", "text/plain")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)


# ---------------------------------------------------------------------------
# Local dev entry-point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    from http.server import HTTPServer

    port = 8000
    print(f"[generate-ppt] Starting -> http://localhost:{port}", flush=True)
    server = HTTPServer(("", port), handler)
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n[generate-ppt] Shutting down.", flush=True)
